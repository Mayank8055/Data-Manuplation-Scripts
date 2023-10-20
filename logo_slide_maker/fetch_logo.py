import pandas as pd
import requests
from io import BytesIO
import os
import shutil
from pptx import Presentation
from pptx.util import Inches

# Load the Excel file
file_path = "Configuration file.xlsm"
sheet_name = "Fetch_logo"
df = pd.read_excel(file_path, sheet_name=sheet_name)

# Specify column names
source_file_col = "Link's Source File"
source_sheet_col = "Link's Source Sheet"
source_column_col = "Link's Source Colunm"
picture_name_col = "Picture Name Colunm"

# Initialize a PowerPoint presentation
ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[6])

# Create a folder to store images
Image_folder = "Downloaded_images"
os.makedirs(Image_folder, exist_ok=True)

# Function to add image to slide
def add_image_to_slide(slide, image_path,name):
    left = Inches(1)
    top = Inches(1)
    width = Inches(6.78)
    height = Inches(7.811)
    pic = slide.shapes.add_picture(image_path, left, top, width, height)

    # Set the name of the picture object
    pic.name = name

# Iterate through the DataFrame
for index, row in df.iterrows():
    source_file = row[source_file_col]
    source_sheet = row[source_sheet_col]
    source_column = row[source_column_col]
    picture_name_col = row[picture_name_col]

    # Read the specified Excel or CSV file
    if source_file.lower().endswith('.csv'):
        excel_df = pd.read_csv(source_file)
    elif source_file.lower().endswith('.xls') or source_file.lower().endswith('.xlsx'):
        excel_df = pd.read_excel(source_file, sheet_name=source_sheet)
    else:
        raise ValueError("Unsupported file format")

    # Process each row in the sub-DataFrame
    for sub_index, sub_row in excel_df.iterrows():
        urls = sub_row[source_column].split(',') if isinstance(sub_row[source_column], str) else ""
        picture_name = sub_row[picture_name_col].strip()

        # Process each URL in the row
        for i, url in enumerate(urls):
            url = url.strip() if isinstance(sub_row[source_column], str) else ""

            # Skip empty URLs
            if url:
                response = requests.get(url)
                if response.status_code == 200:
                    image_data = BytesIO(response.content)
                    image_path = os.path.join(Image_folder, f"{picture_name}_{i}.jpg")
                    img_name = f"{picture_name}_{i}"
                    with open(image_path, "wb") as img_file:
                        img_file.write(image_data.read())
                    add_image_to_slide(slide, image_path,img_name)

# Save the PowerPoint presentation
ppt.save("Images.pptx")

# Delete the folder containing downloaded images
shutil.rmtree(Image_folder)
